# history:
# 2019/03/08  v1.0  initial

import datetime
import os
import msvcrt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

from pylab import mpl
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt


VERSION = '1.0'

START = '1月'
END = '2月'

SAVE_PATH = r'.'
NOTE_PATH = r'./note'
PPT_FILE_NAME_PROFIX = '外协产品线报告'
TITLE = PPT_FILE_NAME_PROFIX
SUB_TITLE = '张文荣'

MARGIN = 0.5

SLD_LAYOUT_TITLE = 0
SLD_LAYOUT_BLANK = 6

XLS_FILE = './budget.xlsx'
BUDGET_AMOUNT = '预算数量'
BUDGET_REVENUE = '预算收入'
BUDGET_PROFIT = '预算毛利'
SALE_AMOUNT = '销售数量'
SALE_REVENUE = '销售收入'
SALE_PROFIT = '销售毛利'

prs = Presentation()
mpl.rcParams['font.sans-serif'] = ['SimHei']
mpl.rcParams['figure.max_open_warning'] = 100
df_budget_amount = pd.read_excel(
    XLS_FILE, sheet_name=BUDGET_AMOUNT, index_col='品名')
df_budget_revenue = pd.read_excel(
    XLS_FILE, sheet_name=BUDGET_REVENUE, index_col='品名')
df_budget_profit = pd.read_excel(
    XLS_FILE, sheet_name=BUDGET_PROFIT, index_col='品名')
df_sale_amount = pd.read_excel(
    XLS_FILE, sheet_name=SALE_AMOUNT, index_col='品名')
df_sale_revenue = pd.read_excel(
    XLS_FILE, sheet_name=SALE_REVENUE, index_col='品名')
df_sale_profit = pd.read_excel(
    XLS_FILE, sheet_name=SALE_PROFIT, index_col='品名')


def info():
    print('=' * 70)
    print('%s v%s' % (os.path.basename(__file__), VERSION))
    print('=' * 70)


def wait_any_key():
    print('=' * 70)
    print('press any key to exit...')
    msvcrt.getch()


def clear_tmp_file():
    os.remove('tmp.png')


def save_ppt():
    date = datetime.datetime.now().strftime('%y%m%d')
    filename = PPT_FILE_NAME_PROFIX + '_' + date + '.pptx'
    print('Save: %s' % filename)
    prs.save(filename)


def set_cell_format(cell):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Pt(2)
    cell.margin_right = Pt(2)
    cell.text_frame.word_wrap = False
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    font = p.font
    font.name = '微软雅黑'
    font.size = Pt(10)


def fill_table(tab, df, digits):
    expr = '%%.%df' % digits

    tab.first_col = True
    # tab.last_col = True
    tab.horz_banding = False
    tab.vert_banding = False

    for i in range(len(df.index)):
        cell = tab.cell(i+1, 0)
        set_cell_format(cell)
        cell.text = df.index[i]

    for i in range(len(df.columns)):
        cell = tab.cell(0, i+1)
        set_cell_format(cell)
        cell.text = df.columns[i]

    for i in range(len(df.index)-1):
        for j in range(len(df.columns)):
            cell = tab.cell(i+1, j+1)
            set_cell_format(cell)
            cell.text = expr % df.iloc[i, j]

    for i in range(len(df.columns)):
        cell = tab.cell(len(df.index), i+1)
        set_cell_format(cell)
        cell.text = '%.0f%%' % (df.iloc[-1, i]*100)


def slide_cover():
    slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_BLANK])

    left = Inches(MARGIN)
    top = Inches(MARGIN)
    width = prs.slide_width - 2*Inches(MARGIN)
    height = prs.slide_height-2*Inches(MARGIN)

    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    font = p.font
    font.name = '华文中宋'
    font.size = Pt(36)
    font.bold = True
    p.text = TITLE
    p.space_after = Pt(20)

    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    font = p.font
    font.name = '华文中宋'
    font.size = Pt(26)
    font.bold = True
    p.text = SUB_TITLE
    print('Slide: cover')


def slide_chart_table(title, tab_df, digits, png_file='tmp.png'):
    slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_BLANK])

    left = Inches(MARGIN)
    top = Inches(MARGIN)
    width = prs.slide_width - 2*Inches(MARGIN)
    height = Inches(MARGIN)

    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    font = p.font
    font.name = '华文中宋'
    font.size = Pt(26)
    font.bold = True
    p.text = title

    top += height
    width = prs.slide_width - 2*Inches(MARGIN)
    height = 9*Inches(MARGIN)
    slide.shapes.add_picture(png_file, left, top, width, height)

    top += height
    width = prs.slide_width - 2*Inches(MARGIN)
    height = 2*Inches(MARGIN)
    rows = tab_df.shape[0]+1
    cols = tab_df.shape[1]+1
    gf = slide.shapes.add_table(rows, cols, left, top, width, height)
    fill_table(gf.table, tab_df, digits)

    print('Slide: %s' % title)


def slide_chart(title, png_file='tmp.png'):
    slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_BLANK])

    left = Inches(MARGIN)
    top = Inches(MARGIN)
    width = prs.slide_width - 2*Inches(MARGIN)
    height = Inches(MARGIN)

    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    font = p.font
    font.name = '华文中宋'
    font.size = Pt(26)
    font.bold = True
    p.text = title

    top += height
    width = prs.slide_width - 2*Inches(MARGIN)
    height = prs.slide_height - height - Inches(MARGIN)
    slide.shapes.add_picture(png_file, left, top, width, height)

    print('Slide: %s' % title)


def slide_note(title, note_file):
    slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_BLANK])

    left = Inches(MARGIN)
    top = Inches(MARGIN)
    width = prs.slide_width - 2*Inches(MARGIN)
    height = Inches(MARGIN)

    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    font = p.font
    font.name = '华文中宋'
    font.size = Pt(26)
    font.bold = True
    p.text = title

    with open(note_file, 'r') as f:
        txt = f.read()

    top += height + Inches(MARGIN)
    width = prs.slide_width - 2*Inches(MARGIN)
    height = prs.slide_height - height - 3*Inches(MARGIN)
    txbox = slide.shapes.add_textbox(left, top, width, height)

    tf = txbox.text_frame
    # tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 2
    font = p.font
    font.name = '华文中宋'
    font.size = Pt(18)
    font.bold = False
    p.text = txt

    print('Slide: %s' % title)


def report_mode1(pd_budget, pd_sale, main_title, sub_title, digit):
    df = pd.DataFrame(index=['预算', '实际', '完成率'], data=[
        pd_budget, pd_sale, pd_sale/pd_budget])
    ax = df.loc[['预算', '实际'], :].T.plot.bar(figsize=(10, 5))
    fig = ax.get_figure()
    fig.savefig('tmp.png')
    df['合计'] = df.apply(lambda x: x.sum(), axis=1)
    df.loc['完成率', '合计'] = df.loc['实际', '合计']/df.loc['预算', '合计']
    slide_chart_table('%s/%s' % (main_title, sub_title), df, digit)


def report_mode2(df, main_title, sub_title):
    ax = df.T.plot.bar(legend=True,
                       stacked=True, figsize=(10, 7))
    fig = ax.get_figure()
    fig.savefig('tmp.png')
    slide_chart('%s/%s' % (main_title, sub_title))


def report_mode3(pd_budget, pd_sale, start, end, sub_title, digit):
    df = pd.DataFrame(index=['预算', '实际', '完成率'], data=[
        pd_budget, pd_sale, pd_sale/pd_budget])
    ax = df.loc[['预算', '实际'], :].T.plot.bar(figsize=(10, 4.5))
    fig = ax.get_figure()
    fig.savefig('tmp.png')
    df['合计'] = df.apply(lambda x: x.sum(), axis=1)
    df.loc['完成率', '合计'] = df.loc['实际', '合计']/df.loc['预算', '合计']
    slide_chart_table('%s-%s/%s' % (start, end, sub_title), df, digit)


def product_report(product):
    a = df_budget_amount.loc[product]
    b = df_sale_amount.loc[product]
    report_mode1(a, b, product, '销售数量', digit=0)
    a = df_budget_revenue.loc[product]
    b = df_sale_revenue.loc[product]
    report_mode1(a, b, product, '销售收入', digit=0)
    a = df_budget_profit.loc[product]
    b = df_sale_profit.loc[product]
    report_mode1(a, b, product, '销售毛利', digit=1)

    note_file = os.path.join(NOTE_PATH, product+'.txt')
    if os.path.exists(note_file):
        slide_note('%s情况说明' % product, note_file)


def main():
    info()

    slide_cover()

    a = df_budget_amount.sum()
    b = df_sale_amount.sum()
    report_mode1(a, b, '总体情况', '总销售数量', digit=0)

    a = df_budget_revenue.sum()
    b = df_sale_revenue.sum()
    report_mode1(a, b, '总体情况', '总销售收入', digit=0)

    a = df_budget_profit.sum()
    b = df_sale_profit.sum()
    report_mode1(a, b, '总体情况', '总销售毛利', digit=1)

    report_mode2(df_sale_revenue, '总体情况', '销售收入构成')
    report_mode2(df_sale_profit, '总体情况', '销售毛利构成')

    a = df_budget_amount.loc[:, START:END].sum(axis=1)
    b = df_sale_amount.loc[:, START:END].sum(axis=1)
    report_mode3(a, b, START, END, '销售数量', digit=0)

    a = df_budget_revenue.loc[:, START:END].sum(axis=1)
    b = df_sale_revenue.loc[:, START:END].sum(axis=1)
    report_mode3(a, b, START, END, '销售收入', digit=0)

    a = df_budget_profit.loc[:, START:END].sum(axis=1)
    b = df_sale_profit.loc[:, START:END].sum(axis=1)
    report_mode3(a, b, START, END, '销售毛利', digit=1)

    product_report('7022')
    product_report('7323')
    product_report('6090')
    product_report('5314')
    product_report('5312')
    product_report('3112')
    # product_report('9006')
    # product_report('HY090')

    save_ppt()
    wait_any_key()
    clear_tmp_file()


if __name__ == '__main__':
    main()
